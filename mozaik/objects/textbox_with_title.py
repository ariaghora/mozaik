from pptx.util import Inches
from pptx.dml.color import RGBColor

from .base_object import BaseObject


class TextboxWithTitle(BaseObject):
    def __init__(self, title, text, title_font_size: float = 0.3):
        self.title = title
        self.text = text
        self.title_font_size = title_font_size

    def attach(self, slide, rect):
        rect.apply_margin()
        textbox = slide.shapes.add_textbox(
            Inches(rect.left_inch),
            Inches(rect.top_inch),
            Inches(rect.width_inch),
            Inches(self.title_font_size + 0.2),
        )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.text = self.title

        # set background color to red
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 0, 0)
        textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        textbox.text_frame.paragraphs[0].font.size = Inches(self.title_font_size)
        textbox.text_frame.paragraphs[0].font.bold = True

        # Add textbox for content
        textbox = slide.shapes.add_textbox(
            Inches(rect.left_inch),
            Inches(rect.top_inch + self.title_font_size + 0.2),
            Inches(rect.width_inch),
            Inches(rect.height_inch - self.title_font_size - 0.2),
        )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.auto_size = False
        textbox.text_frame.text = self.text
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 220, 220)
