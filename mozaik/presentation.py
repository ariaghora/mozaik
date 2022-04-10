"""
TODO: write docstring
"""
from typing import List

from pptx import Presentation as _Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from .picture.picture_modes import set_picture_cover_mode, set_picture_stretch_mode
from .slide import Rect, Slide

config = {
    "slide_title_font_size": 0.5,  # in inches
    "slide_title_font_name": "Arial",
    "slide_title_top_margin": 0.1,  # in inches
    "slide_title_bottom_margin": 0.2,  # in inches
    "slide_title_left_margin": 0.1,  # in inches
    "slide_title_right_margin": 0.1,  # in inches
    "slide_left_padding": 0.1,  # in inches
    "slide_right_padding": 0.1,  # in inches
    "slide_top_padding": 0.1,  # in inches
    "slide_bottom_padding": 0.1,  # in inches
}


class Presentation:
    """
    TODO: write docstring
    """

    def __init__(self, presentation_width: float, presentation_height: float):
        self.slides: List[Slide] = []
        self.presentation_width = presentation_width
        self.presentation_height = presentation_height
        self.prs: _Presentation = None

    def compile_picture(self, slide: Slide, rect: Rect) -> None:
        # Apply margin to the rect
        rect.apply_margin()

        # keep aspect ratio while fitting in the bounding box
        picture_width = Inches(rect.width_inch)
        picture_height = Inches(rect.height_inch)

        # Setup size for fit mode, where the picture is scaled to fit
        # and the ratio is maintained.
        if picture_width < picture_height:
            picture_width = None
        elif picture_width > picture_height:
            picture_height = None

        # Add actual picture. By default, the width and height are
        # set to None, which means that the picture will be scaled
        # to the bounding box proportionally.
        pic = slide.shapes.add_picture(
            rect.content["picture_path"],
            Inches(rect.left_inch),
            Inches(rect.top_inch),
            picture_width,
            picture_height,
        )

        # Set picture mode
        if rect.content["picture_size_mode"] == "fit":
            pass
        elif rect.content["picture_size_mode"] == "stretch":
            set_picture_stretch_mode(pic, rect)
        elif rect.content["picture_size_mode"] == "cover":
            set_picture_cover_mode(pic, rect)
        else:
            raise ValueError(
                f"Unknown picture size mode: {rect.content['picture_size_mode']}"
                + "\nAvailable modes: fit, stretch, cover"
            )

    def compile_text(self, slide: Slide, rect: Rect) -> None:
        rect.apply_margin()
        textbox = slide.shapes.add_textbox(
            Inches(rect.left_inch),
            Inches(rect.top_inch),
            Inches(rect.width_inch),
            Inches(rect.height_inch),
        )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.text = rect.content["text"]

        if rect.content["horizontal_alignment"] == "left":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        elif rect.content["horizontal_alignment"] == "center":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        elif rect.content["horizontal_alignment"] == "right":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        elif rect.content["horizontal_alignment"] == "justify":
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.JUSTIFY
        else:
            raise ValueError(
                f"Unknown horizontal alignment: {rect.content['horizontal_alignment']}"
                + "\nAvailable modes: left, center, right, justify"
            )

    def compile_table(self, slide: Slide, rect: Rect) -> None:
        rect.apply_margin()
        n_rows = len(rect.content["table_data"])
        n_cols = len(rect.content["table_data"][0])

        # This will set table height according to the content
        table_height = Inches(0)

        if rect.content["table_size_mode"] == "stretch":
            table_height = Inches(rect.height_inch)
        elif rect.content["table_size_mode"] == "auto":
            pass
        else:
            raise ValueError(
                f"Unknown table size mode: {rect.content['table_size_mode']}"
                + "\nAvailable modes: stretch, auto"
            )

        table = slide.shapes.add_table(
            n_rows,
            n_cols,
            Inches(rect.left_inch),
            Inches(rect.top_inch),
            Inches(rect.width_inch),
            table_height,
        ).table

        for row_idx, row in enumerate(rect.content["table_data"]):
            for col_idx, cell in enumerate(row):
                table.cell(row_idx, col_idx).text = cell

    def compile_slide_title(self, slide: Slide, title: str) -> None:
        # If the slide has a title, we need to add it as a text box
        textbox = slide.shapes.add_textbox(
            Inches(config["slide_title_left_margin"]),
            Inches(config["slide_title_top_margin"]),
            Inches(
                self.presentation_width
                - config["slide_title_left_margin"]
                - config["slide_title_right_margin"]
            ),
            Inches(config["slide_title_font_size"]),
        )
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.size = Inches(config["slide_title_font_size"])
        paragraph.font.bold = True
        paragraph.text = title

    def compile(self) -> None:
        """
        Compile slides into MS PowerPoint presentation.
        """
        self.prs = _Presentation()
        self.prs.slide_width = Inches(self.presentation_width)
        self.prs.slide_height = Inches(self.presentation_height)

        for slide in self.slides:
            # Add PPT blank layout slide
            __slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

            slide_content_height = self.presentation_height
            if slide.title:
                # shift slide_content_height by title height
                slide_content_height -= (
                    config["slide_title_font_size"]
                    + config["slide_title_top_margin"]
                    + config["slide_title_bottom_margin"]
                )

            # Originally we only have unit sizes (in tiles). We want to convert them
            # to inches. Following method will calculate the size of the slide in inches
            # and keep them in the rect as properties.
            for rect in slide.rects.values():
                rect.calculate_inches(
                    slide,
                    self.presentation_width,
                    slide_content_height,
                    config,
                )

            # If the slide has a title, add it first
            if slide.title:
                self.compile_slide_title(__slide, slide.title)

            # Finally compile all slide contents
            for rect in slide.rects.values():
                if rect.content["type"] == "picture":
                    self.compile_picture(__slide, rect)
                elif rect.content["type"] == "text":
                    self.compile_text(__slide, rect)
                elif rect.content["type"] == "table":
                    self.compile_table(__slide, rect)
                elif rect.content["type"] == "object":
                    rect.content["object"].attach(__slide, rect)

    def add_slide(self, slide: Slide):
        self.slides.append(slide)

    def save(self, name: str):
        self.compile()
        self.prs.save(name)
